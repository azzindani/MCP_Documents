"""Fixtures for the non-PDF formats, built from their own primitives.

Separate from build.py because nothing here is a PDF and the two share no
machinery: build.py writes raw content streams with pikepdf, this writes markup,
zip members and OOXML through the libraries the readers use.

Same rule as build.py, and it is the one that has cost the most to relearn:
**a fixture must not share the property it exists to isolate.** The running-heads
PDF fixture originally templated its body text on the page number, so after
digit-blanking every body line had the same shape as every other -- and the
fixture built to prove headers are stripped instead proved body text could be.

So here: the two-heading HTML file has body paragraphs that are NOT the heading
text repeated; the table fixture's cells are not the same width as its prose;
and the file that tests script-stripping puts a plausible English sentence
inside the <script>, so a stripper that merely drops angle brackets still fails.
"""

from __future__ import annotations

from pathlib import Path

CORPUS = Path(__file__).parent / "corpus"


def _write(name: str, text: str) -> Path:
    CORPUS.mkdir(parents=True, exist_ok=True)
    path = CORPUS / name
    path.write_text(text, encoding="utf-8")
    return path


def article_html(name: str = "article.html") -> Path:
    """A page with declared headings, a table, furniture, and a script.

    The script contains a real sentence rather than obvious code. A reader that
    drops tags without dropping the script SUBTREE leaves it in the text, and
    with `var x = 1` that is easy to spot by eye and easy to miss in an
    assertion; with a sentence, the assertion is the only thing that catches it.
    """
    return _write(
        name,
        """<!doctype html>
<html lang="en-GB">
<head>
  <title>Quarterly Review</title>
  <meta name="description" content="Trading conditions for the period.">
  <style>body { font-family: serif; }</style>
  <script>
    // The board approved the dividend at its meeting in March.
    var total = 41;
  </script>
</head>
<body>
  <header><a href="/">Home</a> &rsaquo; Reports</header>
  <nav><a href="/a">One</a> <a href="/b">Two</a></nav>
  <h1>Quarterly Review</h1>
  <p>Revenue rose across every region we operate in, with the strongest
     contribution from the northern territories.</p>
  <h2>Costs</h2>
  <p>Input prices eased in the second half, though freight remained volatile
     throughout the period under review.</p>
  <table>
    <tr><th>Region</th><th>Revenue</th><th>Change</th></tr>
    <tr><td>North</td><td>412</td><td>+8%</td></tr>
    <tr><td>South</td><td>318</td><td>-2%</td></tr>
  </table>
  <h3>Outlook</h3>
  <ul><li>Continued investment in capacity.</li><li>A review of the dividend.</li></ul>
  <figcaption>Figure 1: revenue by region.</figcaption>
  <footer>Published by the finance team.</footer>
</body>
</html>
""",
    )


def bare_html(name: str = "bare.html") -> Path:
    """Text in a body with no block element at all -- the fallback path.

    Hand-written and generated HTML does this constantly, and a reader that
    only emits blocks for known tags returns an empty document for it.
    """
    return _write(name, "<html><body>Just some words with no paragraph around them at all.</body></html>")


def long_html(name: str = "long.html") -> Path:
    """Long enough to paginate into several synthetic pages.

    Each paragraph names its own index so a test can assert WHICH page a given
    paragraph landed on, rather than only that the count is plausible.
    """
    paras = "\n".join(
        f"  <p>Paragraph {i:03d}. " + ("Filler sentence to give this paragraph some length. " * 6) + "</p>"
        for i in range(1, 61)
    )
    return _write(name, f"<html><body><h1>Long Document</h1>\n{paras}\n</body></html>")


def readme_md(name: str = "readme.md") -> Path:
    """Markdown with ATX and setext headings, and a fenced block that lies.

    The fenced block contains `# Not a heading`, which is a shell comment. A
    reader that scans every line for a leading hash puts it in the outline.
    """
    return _write(
        name,
        """# Project Title

An opening paragraph describing what this is for.

Second Level Heading
--------------------

Some prose under the setext heading.

## Installation

```sh
# Not a heading, this is a shell comment
make install
```

### Notes

Closing remarks about the project.
""",
    )


def sales_csv(name: str = "sales.csv") -> Path:
    """A CSV whose quoted fields contain the delimiter and a newline.

    Splitting on commas produces more columns on some rows than others, with
    the values shifted -- data that is wrong rather than obviously missing,
    which is the failure this fixture exists to catch.
    """
    return _write(
        name,
        "region,note,revenue\n"
        'North,"Strong, steady growth",412\n'
        'South,"Freight costs rose;\nmargins narrowed",318\n'
        "East,Flat,207\n",
    )


def plain_txt(name: str = "notes.txt") -> Path:
    return _write(
        name,
        "First paragraph of a plain text file.\nStill the first paragraph.\n"
        "\n"
        "A second paragraph, separated by a blank line.\n"
        "\n"
        "A third one to make the block count worth asserting.\n",
    )


def message_eml(name: str = "message.eml") -> Path:
    """A multipart message with both bodies and an attachment.

    The HTML alternative says something the plain part does not, so a reader
    that returns the wrong one -- or both -- is caught by content rather than
    by length.
    """
    return _write(
        name,
        """From: Alex Fenner <alex@example.com>
To: Sam Reyes <sam@example.com>
Subject: =?utf-8?B?UXVhcnRlcmx5IHJldmlldyBhdHRhY2hlZA==?=
Date: Mon, 4 Aug 2026 09:14:02 +0100
MIME-Version: 1.0
Content-Type: multipart/mixed; boundary="OUTER"

--OUTER
Content-Type: multipart/alternative; boundary="INNER"

--INNER
Content-Type: text/plain; charset="utf-8"

Please find the quarterly review attached.

The board meets on Thursday.
--INNER
Content-Type: text/html; charset="utf-8"

<html><body><p>HTML ALTERNATIVE BODY</p></body></html>
--INNER--
--OUTER
Content-Type: text/csv; name="figures.csv"
Content-Disposition: attachment; filename="figures.csv"

region,revenue
North,412
--OUTER--
""",
    )


def html_only_eml(name: str = "html_only.eml") -> Path:
    """A message with no plain-text part, so the HTML path must be used."""
    return _write(
        name,
        """From: noreply@example.com
To: sam@example.com
Subject: Statement ready
Date: Tue, 5 Aug 2026 08:00:00 +0100
MIME-Version: 1.0
Content-Type: text/html; charset="utf-8"

<html><body><script>var t = 1;</script>
<h1>Statement ready</h1><p>Your statement for July is available.</p></body></html>
""",
    )


def report_docx(name: str = "report.docx") -> Path:
    """A Word document with a table BETWEEN two paragraphs.

    Position is the whole point. python-docx exposes paragraphs and tables as
    two flat lists, and reading them in that order puts every table at the end
    of the document -- a plausible, readable, wrongly-ordered result.
    """
    import docx

    document = docx.Document()
    document.add_heading("Annual Report", level=1)
    document.add_paragraph("Opening paragraph that comes before the table.")
    document.add_heading("Figures", level=2)
    table = document.add_table(rows=3, cols=2)
    for row, (left, right) in enumerate((("Region", "Revenue"), ("North", "412"), ("South", "318"))):
        table.rows[row].cells[0].text = left
        table.rows[row].cells[1].text = right
    document.add_paragraph("Closing paragraph that comes after the table.")

    CORPUS.mkdir(parents=True, exist_ok=True)
    path = CORPUS / name
    document.save(str(path))
    return path


def deck_pptx(name: str = "deck.pptx") -> Path:
    """Three slides, one with speaker notes, one with a table.

    The body text box is added BEFORE the title on slide 2 so shape order and
    visual order disagree -- a reader sorting by shape index reads the body
    first.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    blank = prs.slide_layouts[6]
    titled = prs.slide_layouts[5]

    first = prs.slides.add_slide(titled)
    title = first.shapes.title
    assert title is not None, "the title layout should carry a title placeholder"
    title.text = "Opening Slide"
    notes = first.notes_slide.notes_text_frame
    assert notes is not None, "a notes slide always carries a notes text frame"
    notes.text = "Speaker note: mention the revised timetable."

    second = prs.slides.add_slide(blank)
    body = second.shapes.add_textbox(Inches(1), Inches(3), Inches(6), Inches(1))
    body.text_frame.text = "Body text placed lower on the slide."
    heading = second.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    heading.text_frame.text = "Second Slide Heading"
    heading.text_frame.paragraphs[0].runs[0].font.size = Pt(32)

    third = prs.slides.add_slide(blank)
    shape = third.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(1)).table
    shape.cell(0, 0).text = "Region"
    shape.cell(0, 1).text = "Revenue"
    shape.cell(1, 0).text = "North"
    shape.cell(1, 1).text = "412"

    CORPUS.mkdir(parents=True, exist_ok=True)
    path = CORPUS / name
    prs.save(str(path))
    return path


def book_xlsx(name: str = "book.xlsx") -> Path:
    """Two sheets, a formula with a cached value, and a far-flung stray cell.

    The stray cell at Z400 makes `max_row` 400 while there are four real rows,
    so a reader trusting the sheet dimensions reports a document 100x its real
    size.
    """
    import openpyxl

    book = openpyxl.Workbook()
    first = book.active
    assert first is not None, "a new workbook always has an active sheet"
    first.title = "Revenue"
    for row in (("Region", "Revenue"), ("North", 412), ("South", 318), ("Total", "=SUM(B2:B3)")):
        first.append(row)
    # A number format on a far cell is what makes max_row 400 for a four-row
    # sheet: the point of this fixture is a sheet whose DIMENSIONS lie, which
    # in real exports comes from a format applied to a whole column.
    #
    # Measured, because the obvious way does not work: `first["Z400"] = None`
    # leaves max_row at 4 after a save/reload, so the fixture did not have the
    # property it exists to isolate and the test guarding it passed for the
    # wrong reason. `= ""`, `= 0` and a bare number_format all give 400.
    first["Z400"].number_format = "0.00"

    second = book.create_sheet("Notes")
    second.append(("Note",))
    second.append(("Freight costs rose in the second half.",))

    CORPUS.mkdir(parents=True, exist_ok=True)
    path = CORPUS / name
    book.save(str(path))
    return path


def book_epub(name: str = "book.epub") -> Path:
    """A minimal but real epub whose SPINE order is not its filename order.

    Chapters are named c3, c1, c2 in the archive and ordered 1, 2, 3 in the
    spine. A reader using the manifest, the archive order, or an alphabetical
    sort gets a plausible book with its chapters shuffled.
    """
    import zipfile

    CORPUS.mkdir(parents=True, exist_ok=True)
    path = CORPUS / name

    def chapter(number: int, title: str, body: str) -> str:
        return (
            '<?xml version="1.0" encoding="utf-8"?>\n'
            '<html xmlns="http://www.w3.org/1999/xhtml"><head><title>'
            f"{title}</title></head><body><h1>{title}</h1><p>{body}</p></body></html>"
        )

    container = (
        '<?xml version="1.0"?>\n'
        '<container version="1.0" xmlns="urn:oasis:names:tc:opendocument:xmlns:container">'
        '<rootfiles><rootfile full-path="OEBPS/book.opf" media-type="application/oebps-package+xml"/>'
        "</rootfiles></container>"
    )
    opf = (
        '<?xml version="1.0"?>\n'
        '<package xmlns="http://www.idpf.org/2007/opf" version="3.0" unique-identifier="id">'
        '<metadata xmlns:dc="http://purl.org/dc/elements/1.1/">'
        '<dc:identifier id="id">urn:uuid:test</dc:identifier>'
        "<dc:title>A Short Book</dc:title>"
        "<dc:creator>R. Fenner</dc:creator>"
        "<dc:language>en</dc:language>"
        "</metadata><manifest>"
        '<item id="c3" href="c3.xhtml" media-type="application/xhtml+xml"/>'
        '<item id="c1" href="c1.xhtml" media-type="application/xhtml+xml"/>'
        '<item id="c2" href="c2.xhtml" media-type="application/xhtml+xml"/>'
        "</manifest><spine>"
        '<itemref idref="c1"/><itemref idref="c2"/><itemref idref="c3"/>'
        "</spine></package>"
    )

    with zipfile.ZipFile(path, "w") as archive:
        # mimetype first and STORED, as the spec requires.
        archive.writestr("mimetype", "application/epub+zip", compress_type=zipfile.ZIP_STORED)
        archive.writestr("META-INF/container.xml", container)
        archive.writestr("OEBPS/book.opf", opf)
        archive.writestr("OEBPS/c3.xhtml", chapter(3, "Chapter Three", "The last chapter of the book."))
        archive.writestr("OEBPS/c1.xhtml", chapter(1, "Chapter One", "The first chapter of the book."))
        archive.writestr("OEBPS/c2.xhtml", chapter(2, "Chapter Two", "The middle chapter of the book."))
    return path


BUILDERS = {
    "article_html": article_html,
    "bare_html": bare_html,
    "long_html": long_html,
    "readme_md": readme_md,
    "sales_csv": sales_csv,
    "plain_txt": plain_txt,
    "message_eml": message_eml,
    "html_only_eml": html_only_eml,
    "report_docx": report_docx,
    "deck_pptx": deck_pptx,
    "book_xlsx": book_xlsx,
    "book_epub": book_epub,
}


def build_all() -> dict[str, Path]:
    return {name: fn() for name, fn in BUILDERS.items()}


if __name__ == "__main__":
    for name, built in build_all().items():
        print(f"{name:<16} {built.stat().st_size:>9,} bytes  {built}")
