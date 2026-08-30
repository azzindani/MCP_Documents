"""PowerPoint -> core.ir.Document, via python-pptx.

**One slide is one page, and that is a real boundary**, not a synthetic one.
The file says where each slide begins; nothing has to be inferred and nothing
is disclosed as approximate. This is the difference `_flow.paged` exists to
keep: telling a caller a deck's pages are synthetic would be as wrong as
telling them an HTML file's are real.

Slide order is presentation order. python-pptx's `prs.slides` follows the
sldIdLst in the presentation part, which is the order they are shown in --
not the order they happen to sit in the package.

Speaker notes are read and marked `head_foot`, so the cleaner can drop them and
a caller who wants them can still see them. They are frequently where the
actual content of a deck lives, and dropping them at read time would make that
unrecoverable.
"""

from __future__ import annotations

from core.ir import Block, Document, Page, Span
from core.readers import _flow

TITLE_SIZE = 26.0
BODY_SIZE = 12.0


def open_document(path: str, password: str = "") -> Document:
    """Read a .pptx. One page per slide."""
    from pptx import Presentation

    src = _flow.guard_size(path)
    try:
        handle = Presentation(str(src))
    except Exception as exc:
        from core.readers import ReaderError

        raise ReaderError(
            f"{src.name} could not be read as a PowerPoint file: {exc}",
            "Check the file opens in PowerPoint. A .ppt (not .pptx) needs convert(to='pdf') first.",
        ) from exc

    width = float(handle.slide_width or 0) / 12700.0  # EMU -> points
    height = float(handle.slide_height or 0) / 12700.0

    pages: list[Page] = []
    tables = 0
    for index, slide in enumerate(handle.slides, start=1):
        blocks = _slide_blocks(slide)
        tables += sum(1 for b in blocks if b.kind == "table")
        page = Page(number=index, width=width, height=height, basis="native" if blocks else "empty")
        for order, block in enumerate(blocks):
            block.order = order
        page.blocks = blocks
        pages.append(page)

    core = handle.core_properties
    meta = {
        "title": core.title or "",
        "author": core.author or "",
        "slides": len(pages),
        "tables": tables,
    }
    return _flow.paged(str(src), "pptx", pages, meta)


def _slide_blocks(slide) -> list[Block]:
    """Every shape's text, in the order the slide places them.

    Sorted by position -- top then left -- rather than by shape index. Shape
    order in the XML is z-order, which is the order they were ADDED, so a title
    someone moved to the front of the stack after writing the body would read
    last. Slides really do have geometry, so this is a measurement, not a
    guess.
    """
    placed: list[tuple[float, float, Block]] = []

    for shape in slide.shapes:
        top = float(shape.top or 0) / 12700.0
        left = float(shape.left or 0) / 12700.0

        if getattr(shape, "has_table", False):
            placed.append((top, left, _table(shape.table)))
            continue
        if not getattr(shape, "has_text_frame", False):
            continue
        text = "\n".join(p.text for p in shape.text_frame.paragraphs if p.text.strip()).strip()
        if not text:
            continue

        is_title = shape == slide.shapes.title
        kind = "heading" if is_title else "para"
        size = TITLE_SIZE if is_title else BODY_SIZE
        block = _flow.block(text, kind, "native", size)
        # Slides have real coordinates, so blocks carry them: reading order,
        # column detection and redaction all work on a deck the same way they
        # work on a PDF, rather than falling back to "keep the reader's order".
        block.bbox = (
            left,
            top,
            left + float(shape.width or 0) / 12700.0,
            top + float(shape.height or 0) / 12700.0,
        )
        for span in block.spans:
            span.bbox = block.bbox
        placed.append((top, left, block))

    blocks = [block for _, _, block in sorted(placed, key=lambda item: (item[0], item[1]))]

    if slide.has_notes_slide:
        notes = (slide.notes_slide.notes_text_frame.text or "").strip()
        if notes:
            blocks.append(_flow.block(notes, "head_foot", "native", BODY_SIZE))
    return blocks


def _table(table) -> Block:
    rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
    rows = [row for row in rows if any(cell for cell in row)]
    flat = "\n".join("\t".join(cell for cell in row) for row in rows)
    return Block(kind="table", spans=[Span(text=flat, size=BODY_SIZE)], bbox=None, basis="native", rows=rows)


def page_tables(doc: Document, numbers: list[int]) -> list[dict]:
    from core.readers.html import page_tables as _shared

    return _shared(doc, numbers)


def bookmarks(doc: Document) -> list[dict]:
    """One entry per slide title -- a deck's outline is its titles."""
    out: list[dict] = []
    for number in sorted(doc.pages):
        for block in doc.pages[number].blocks:
            if block.kind == "heading":
                out.append({"level": 1, "title": block.text.strip(), "page": number, "basis": "native"})
                break
    return out


def probe_extras(doc: Document) -> dict:
    return {key: doc.meta[key] for key in ("title", "author", "slides", "tables") if key in doc.meta}


close_document = _flow.close_document
load_page = _flow.load_page
load_page_words = _flow.load_page
